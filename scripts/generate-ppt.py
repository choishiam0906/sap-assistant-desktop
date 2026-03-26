"""
SAP Knowledge Hub v7.1.0 — 경영진/투자자 발표 PPT 생성
python-pptx를 사용하여 .pptx 파일 생성
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 색상 팔레트
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)       # 진한 네이비
ACCENT_BLUE = RGBColor(0x38, 0x7A, 0xFF)    # 포인트 블루
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)   # 포인트 그린
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)  # 포인트 오렌지
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)     # 포인트 레드
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
MID_GRAY = RGBColor(0x64, 0x74, 0x8B)
CARD_BG = RGBColor(0x1E, 0x29, 0x3B)        # 카드 배경

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=DARK_BG):
    """슬라이드 배경색 설정"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, alignment=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_shape(slide, left, top, width, height, color=CARD_BG):
    """둥근 사각형 카드 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_card_with_text(slide, left, top, width, height, title, body,
                        accent=ACCENT_BLUE):
    """카드 + 텍스트 조합"""
    add_shape(slide, left, top, width, height)
    add_text(slide, left + 0.3, top + 0.2, width - 0.6, 0.4,
             title, size=16, color=accent, bold=True)
    add_text(slide, left + 0.3, top + 0.7, width - 0.6, height - 0.9,
             body, size=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════
# Slide 1: 타이틀
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# 상단 태그
add_text(slide, 1, 1.2, 11, 0.5, "PRODUCT PRESENTATION  ·  v7.1.0  ·  2026.03",
         size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# 메인 타이틀
add_text(slide, 1, 2.2, 11, 1.2, "SAP Knowledge Hub",
         size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 서브 타이틀
add_text(slide, 1.5, 3.6, 10, 0.8,
         "제조기업 현업 운영자를 위한 AI 기반 에러 자가 진단 플랫폼",
         size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# 핵심 한 줄
add_text(slide, 2, 4.8, 9, 0.6,
         "코드 · 이메일 · 문서를 통합 분석하여 시스템 에러를 스스로 해결하세요",
         size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# 하단 태그
add_text(slide, 1, 6.2, 11, 0.4,
         "Local-First  ·  Multi-LLM  ·  18 Languages  ·  Zero Data Leak",
         size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# Slide 2: 문제 정의
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 0.8, 0.5, 11, 0.6, "문제 정의", size=32, color=WHITE, bold=True)
add_text(slide, 0.8, 1.1, 11, 0.4,
         "제조기업 현업 운영자가 직면한 5가지 핵심 고통점",
         size=16, color=MID_GRAY)

problems = [
    ("🔴  에러 이해 불가", "SAP/MES/QMS 에러 메시지가\n비전문가에게는 암호화 수준"),
    ("🟠  IT 병목 현상", "단순 질문도 IT팀 대기 평균 4시간\n→ 업무 중단, 생산성 저하"),
    ("🟡  반복 문제 재발", "동일 에러가 반복되지만\n해결 지식이 공유되지 않음"),
    ("🔵  문서 접근성 부재", "SAP Note, 코드, 매뉴얼이\n분산 저장 → 찾기 어려움"),
    ("🟣  이메일 과부하", "운영 관련 이메일이 분산되어\n중요 알림을 놓치는 문제"),
]

for i, (title, body) in enumerate(problems):
    left = 0.8 + (i % 5) * 2.4
    top = 2.0
    add_card_with_text(slide, left, top, 2.2, 3.2, title, body, ACCENT_ORANGE)

# 하단 강조
add_shape(slide, 0.8, 5.8, 11.7, 0.8, RGBColor(0x1A, 0x25, 0x3A))
add_text(slide, 1.2, 5.9, 11, 0.6,
         "💡 결과: IT 티켓의 60%가 \"이미 해결된 문제\"의 반복 — AI로 자가 진단하면 즉시 해소 가능",
         size=15, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# Slide 3: 솔루션 개요
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 0.8, 0.5, 11, 0.6, "솔루션: SAP Knowledge Hub", size=32, color=WHITE, bold=True)
add_text(slide, 0.8, 1.1, 11, 0.4,
         "3가지 데이터 소스를 AI로 통합 분석하는 로컬 우선 플랫폼",
         size=16, color=MID_GRAY)

# 3개 핵심 기둥
pillars = [
    ("📂  코드 분석", "GitHub 리포지토리\nREST API 연동\n18개 언어 지원\nSHA 증분 동기화", ACCENT_BLUE),
    ("📧  이메일 분석", "Gmail + Outlook 연동\nSAP 메타데이터 추출\nAI 분석 → 플랜 자동 생성\n수동 임포트 지원", ACCENT_GREEN),
    ("💬  다중 LLM 채팅", "OpenAI GPT-4.1\nClaude Opus/Sonnet\nGemini 2.5\nSSE 실시간 스트리밍", ACCENT_ORANGE),
]

for i, (title, body, color) in enumerate(pillars):
    left = 0.8 + i * 4.1
    add_shape(slide, left, 2.0, 3.8, 3.0)
    add_text(slide, left + 0.3, 2.2, 3.2, 0.5, title, size=20, color=color, bold=True)
    add_text(slide, left + 0.3, 2.8, 3.2, 2.0, body, size=14, color=LIGHT_GRAY)

# 하단: 보안 강조
add_shape(slide, 0.8, 5.5, 11.7, 1.2, RGBColor(0x1A, 0x25, 0x3A))
add_text(slide, 1.2, 5.6, 2.5, 0.4, "🔒 보안 아키텍처", size=16, color=WHITE, bold=True)

security_items = [
    "Local-First: 데이터가 사용자 PC를 떠나지 않음",
    "AES-256-GCM 암호화 + 시스템 키체인(keytar)",
    "OAuth 2.0 + PKCE (4대 프로바이더)",
    "API 키 자동 마스킹 (Logger Redaction)",
]
for i, item in enumerate(security_items):
    add_text(slide, 1.2 + (i % 2) * 5.5, 6.05 + (i // 2) * 0.3, 5.5, 0.3,
             f"✓  {item}", size=12, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════
# Slide 4: 제품 데모 — 핵심 화면
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 0.8, 0.5, 11, 0.6, "제품 구조", size=32, color=WHITE, bold=True)
add_text(slide, 0.8, 1.1, 11, 0.4,
         "Electron 데스크톱 앱 — 7개 주요 화면",
         size=16, color=MID_GRAY)

screens = [
    ("Chat", "다중 LLM\n실시간 대화", ACCENT_BLUE),
    ("CodeLab", "코드 분석\nGitHub 연동", ACCENT_GREEN),
    ("Cockpit", "마감 관리\n스케줄 자동화", ACCENT_ORANGE),
    ("Knowledge", "에이전트\n스킬 카탈로그", RGBColor(0xA7, 0x8B, 0xFA)),
    ("Email", "이메일 동기화\nAI 분석", RGBColor(0xF4, 0x72, 0xB6)),
    ("Settings", "AI/OAuth\nCodeLab 설정", MID_GRAY),
    ("Audit", "감사 로그\n세션 이력", RGBColor(0x6E, 0xE7, 0xB7)),
]

for i, (name, desc, color) in enumerate(screens):
    left = 0.5 + i * 1.8
    add_shape(slide, left, 2.0, 1.6, 2.4)
    add_text(slide, left + 0.1, 2.15, 1.4, 0.4, name, size=15, color=color, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(slide, left + 0.1, 2.6, 1.4, 1.5, desc, size=12, color=LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER)

# 기술 스택 바
add_shape(slide, 0.8, 4.8, 11.7, 1.8, RGBColor(0x1A, 0x25, 0x3A))
add_text(slide, 1.2, 4.9, 4, 0.4, "기술 스택", size=16, color=WHITE, bold=True)

stack_items = [
    ("Runtime", "Electron 31 + Node.js 22"),
    ("Frontend", "React 18 + TypeScript 5.7"),
    ("State", "Zustand 5 + React Query 5"),
    ("Database", "SQLite (better-sqlite3)"),
    ("Auth", "OAuth 2.0 PKCE + keytar"),
    ("LLM", "OpenAI / Anthropic / Google"),
    ("Protocol", "MCP SDK 1.27"),
    ("Build", "Vite 6 + esbuild"),
]

for i, (k, v) in enumerate(stack_items):
    col = i % 4
    row = i // 4
    add_text(slide, 1.2 + col * 2.9, 5.4 + row * 0.45, 1.0, 0.3,
             k, size=11, color=ACCENT_BLUE, bold=True)
    add_text(slide, 2.2 + col * 2.9, 5.4 + row * 0.45, 2.0, 0.3,
             v, size=11, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════
# Slide 5: 시장 기회 & 타겟
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 0.8, 0.5, 11, 0.6, "시장 기회", size=32, color=WHITE, bold=True)

# TAM/SAM/SOM
add_card_with_text(slide, 0.8, 1.5, 3.7, 2.5,
    "TAM — 전체 시장",
    "글로벌 ERP 시장\n$62.5B (2025)\nAI 기반 운영 도구\n연 18% 성장률",
    ACCENT_BLUE)

add_card_with_text(slide, 4.8, 1.5, 3.7, 2.5,
    "SAM — 접근 가능 시장",
    "한국 + 아시아 제조기업\nSAP/MES/QMS 사용 기업\n약 5,000개사\n$2.1B",
    ACCENT_GREEN)

add_card_with_text(slide, 8.8, 1.5, 3.7, 2.5,
    "SOM — 목표 시장",
    "대규모 제조기업 (1000인+)\nSAP ERP 운영 중\n초기 타겟 500개사\n$150M",
    ACCENT_ORANGE)

# 타겟 페르소나
add_text(slide, 0.8, 4.3, 6, 0.4, "타겟 사용자 페르소나", size=18, color=WHITE, bold=True)

personas = [
    ("👷 현업 운영자", "시스템 에러 시 스스로\n해결책을 찾고 싶은 사람\n→ 핵심 타겟 (80%)"),
    ("👨‍💻 IT 관리자", "반복 티켓을 줄이고\n운영 효율화를 원하는 사람\n→ 구매 결정자"),
    ("🏭 공장장/팀장", "생산 중단 시간을\n최소화하고 싶은 사람\n→ 예산 승인자"),
]

for i, (title, desc) in enumerate(personas):
    add_card_with_text(slide, 0.8 + i * 4.1, 4.8, 3.8, 2.2, title, desc, ACCENT_BLUE)


# ═══════════════════════════════════════════════
# Slide 6: ROI & 비즈니스 가치
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 0.8, 0.5, 11, 0.6, "ROI 분석", size=32, color=WHITE, bold=True)
add_text(slide, 0.8, 1.1, 11, 0.4,
         "투자 대비 수익률 — 도입 4개월 만에 손익분기점 돌파",
         size=16, color=MID_GRAY)

# 핵심 수치 카드
metrics = [
    ("267%", "Year 1 ROI", ACCENT_GREEN),
    ("4개월", "손익분기점", ACCENT_BLUE),
    ("40%↓", "IT 티켓 감소", ACCENT_ORANGE),
    ("80%↓", "MTTR 단축", RGBColor(0xA7, 0x8B, 0xFA)),
]

for i, (value, label, color) in enumerate(metrics):
    left = 0.8 + i * 3.1
    add_shape(slide, left, 1.8, 2.8, 1.8)
    add_text(slide, left, 2.0, 2.8, 0.8, value, size=40, color=color, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(slide, left, 2.9, 2.8, 0.5, label, size=16, color=LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER)

# 비용/절감 테이블
add_shape(slide, 0.8, 4.0, 5.5, 3.0, RGBColor(0x1A, 0x25, 0x3A))
add_text(slide, 1.1, 4.1, 4, 0.4, "💰 투자 비용", size=16, color=WHITE, bold=True)
costs = [
    "개발 인건비 (6개월):     80M KRW",
    "인프라 (클라우드/LLM):    30M KRW",
    "라이선스 (Electron 등):   20M KRW",
    "마케팅/영업:              20M KRW",
    "──────────────────────",
    "총 초기 투자:            150M KRW",
]
for i, line in enumerate(costs):
    c = ACCENT_GREEN if "총" in line else LIGHT_GRAY
    add_text(slide, 1.1, 4.6 + i * 0.35, 5, 0.3, line, size=12, color=c,
             font_name="Consolas")

add_shape(slide, 7.0, 4.0, 5.5, 3.0, RGBColor(0x1A, 0x25, 0x3A))
add_text(slide, 7.3, 4.1, 4, 0.4, "📈 연간 절감 효과", size=16, color=WHITE, bold=True)
savings = [
    "IT 지원 인건비 절감:     250M KRW",
    "생산 중단 시간 감소:     200M KRW",
    "반복 에러 해소:          100M KRW",
    "문서/이메일 자동화:       50M KRW",
    "──────────────────────",
    "총 연간 절감:            600M KRW",
]
for i, line in enumerate(savings):
    c = ACCENT_GREEN if "총" in line else LIGHT_GRAY
    add_text(slide, 7.3, 4.6 + i * 0.35, 5, 0.3, line, size=12, color=c,
             font_name="Consolas")


# ═══════════════════════════════════════════════
# Slide 7: 경쟁 우위
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 0.8, 0.5, 11, 0.6, "경쟁 우위", size=32, color=WHITE, bold=True)
add_text(slide, 0.8, 1.1, 11, 0.4,
         "SAP Knowledge Hub vs 기존 솔루션",
         size=16, color=MID_GRAY)

# 비교 테이블 헤더
headers = ["기능", "SAP Knowledge Hub", "ChatGPT/Copilot", "ServiceNow", "SAP Solution Mgr"]
col_widths = [2.2, 2.5, 2.5, 2.3, 2.5]
x_start = 0.5

for i, (h, w) in enumerate(zip(headers, col_widths)):
    x = x_start + sum(col_widths[:i])
    bg_color = ACCENT_BLUE if i == 1 else CARD_BG
    add_shape(slide, x, 1.8, w - 0.1, 0.5, bg_color)
    add_text(slide, x + 0.1, 1.85, w - 0.3, 0.4, h, size=12, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

rows = [
    ["로컬 우선 보안", "✅ 완전 로컬", "❌ 클라우드", "❌ 클라우드", "△ 온프레미스"],
    ["다중 LLM", "✅ 3사 통합", "❌ GPT 전용", "△ 제한적", "❌ 미지원"],
    ["코드 분석", "✅ GitHub 연동", "△ 수동 붙여넣기", "❌ 미지원", "△ 제한적"],
    ["이메일 분석", "✅ AI 자동 분석", "❌ 미지원", "△ 규칙 기반", "❌ 미지원"],
    ["SAP 특화", "✅ 도메인 팩 5종", "❌ 범용", "△ ITSM 중심", "✅ SAP 전용"],
    ["비용", "✅ 라이선스 무료", "월 $20/인", "월 $100+/인", "높은 유지비"],
    ["커스터마이징", "✅ agent.md/skill.md", "❌ 제한적", "△ 워크플로우", "△ 복잡"],
]

for r, row in enumerate(rows):
    y = 2.4 + r * 0.55
    for c, (cell, w) in enumerate(zip(row, col_widths)):
        x = x_start + sum(col_widths[:c])
        stripe = RGBColor(0x16, 0x20, 0x33) if r % 2 == 0 else CARD_BG
        if c == 1:
            stripe = RGBColor(0x1A, 0x2E, 0x4A) if r % 2 == 0 else RGBColor(0x15, 0x26, 0x3F)
        add_shape(slide, x, y, w - 0.1, 0.5, stripe)
        color = ACCENT_GREEN if cell.startswith("✅") else (ACCENT_RED if cell.startswith("❌") else LIGHT_GRAY)
        add_text(slide, x + 0.1, y + 0.05, w - 0.3, 0.4, cell, size=11, color=color,
                 alignment=PP_ALIGN.CENTER)

# 핵심 차별점
add_shape(slide, 0.8, 6.4, 11.7, 0.7, RGBColor(0x1A, 0x25, 0x3A))
add_text(slide, 1.2, 6.5, 11, 0.5,
         "핵심 차별점:  데이터 유출 0%  ·  다중 LLM 선택  ·  코드+이메일+문서 통합 분석  ·  커스텀 에이전트",
         size=14, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# Slide 8: 로드맵
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 0.8, 0.5, 11, 0.6, "로드맵", size=32, color=WHITE, bold=True)
add_text(slide, 0.8, 1.1, 11, 0.4,
         "v4.0 → v7.1 (완료) → v8.0+ (계획)",
         size=16, color=MID_GRAY)

phases = [
    ("v4.0", "수동 워크플로우", "2025 Q4", "다중 LLM 채팅\n에이전트/스킬 시스템\nCBO 분석\nCode Lab", ACCENT_BLUE, True),
    ("v5.0", "자동화", "2026 Q1", "SSE 스트리밍\n스케줄 자동화\n에러 복원력\nDB 마이그레이션", ACCENT_BLUE, True),
    ("v6.x", "플랫폼 성숙", "2026 Q1", "UI 분할/a11y\n범용 플랫폼 전환\nReact Query 최적화\nZustand persist", ACCENT_BLUE, True),
    ("v7.1", "Knowledge Hub", "2026 Q1", "GitHub 코드 연동\n이메일 AI 분석\nOAuth 4 프로바이더\n보안 강화", ACCENT_GREEN, True),
    ("v8.0", "엔터프라이즈", "2026 H2", "GitLab 연동\n오프라인 LLM\nLDAP/SSO\n코드 서명", ACCENT_ORANGE, False),
]

for i, (ver, name, date, features, color, done) in enumerate(phases):
    left = 0.3 + i * 2.55
    # 타임라인 연결선
    if i < len(phases) - 1:
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left + 2.1), Inches(2.3), Inches(0.5), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = MID_GRAY
        line.line.fill.background()

    # 버전 원
    circle_color = color if done else MID_GRAY
    add_shape(slide, left + 0.7, 1.8, 1.0, 1.0, circle_color)
    add_text(slide, left + 0.7, 2.0, 1.0, 0.5, ver, size=16, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

    # 이름 + 날짜
    add_text(slide, left + 0.1, 3.0, 2.2, 0.4, name, size=14, color=color, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(slide, left + 0.1, 3.35, 2.2, 0.3, date, size=11, color=MID_GRAY,
             alignment=PP_ALIGN.CENTER)

    # 기능 카드
    status = "✅ 완료" if done else "🔮 계획"
    add_shape(slide, left + 0.1, 3.8, 2.2, 2.8)
    add_text(slide, left + 0.2, 3.9, 2.0, 0.3, status, size=11, color=color, bold=True)
    add_text(slide, left + 0.2, 4.3, 2.0, 2.2, features, size=12, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════
# Slide 9: 현재 성과 & KPI
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 0.8, 0.5, 11, 0.6, "현재 성과 (v7.1.0)", size=32, color=WHITE, bold=True)

# 기술 KPI
tech_metrics = [
    ("167+", "IPC 채널", "Electron Main ↔ Renderer"),
    ("108", "Preload 메서드", "window.assistantDesktop API"),
    ("277", "테스트 통과", "Vitest + RTL (19 신규)"),
    ("20+", "Repository", "SQLite CRUD 추상화"),
]

for i, (val, label, desc) in enumerate(tech_metrics):
    left = 0.8 + i * 3.1
    add_shape(slide, left, 1.3, 2.8, 1.6)
    add_text(slide, left, 1.4, 2.8, 0.7, val, size=36, color=ACCENT_BLUE, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(slide, left, 2.1, 2.8, 0.35, label, size=14, color=WHITE,
             alignment=PP_ALIGN.CENTER)
    add_text(slide, left, 2.4, 2.8, 0.35, desc, size=11, color=MID_GRAY,
             alignment=PP_ALIGN.CENTER)

# 기능 체크리스트
add_text(slide, 0.8, 3.2, 6, 0.4, "구현 완료 기능 체크리스트", size=18, color=WHITE, bold=True)

features_left = [
    "✅ 다중 LLM 채팅 (OpenAI/Anthropic/Google)",
    "✅ SSE 실시간 스트리밍",
    "✅ CBO 분석 (규칙 + LLM 강화)",
    "✅ 커스텀 에이전트/스킬 (agent.md)",
    "✅ 5가지 도메인 팩 (Ops, Func, CBO, PI, BTP)",
    "✅ Knowledge Vault + MCP 서버",
    "✅ Cockpit 대시보드 (마감 + 스케줄)",
    "✅ DB 마이그레이션 시스템",
]

features_right = [
    "✅ OAuth 4대 프로바이더 (PKCE)",
    "✅ GitHub 코드 연동 (REST API)",
    "✅ 이메일 통합 분석 (Gmail/Outlook)",
    "✅ 코드 분석 엔진 (18개 언어)",
    "✅ SecureStore (keytar + AES-256)",
    "✅ Logger Redaction (API 키 마스킹)",
    "✅ 에러 복원력 (Retry + CB + Fallback)",
    "✅ NSIS + Portable EXE 배포",
]

for i, item in enumerate(features_left):
    add_text(slide, 0.8, 3.7 + i * 0.4, 6, 0.35, item, size=12, color=LIGHT_GRAY)

for i, item in enumerate(features_right):
    add_text(slide, 7.0, 3.7 + i * 0.4, 6, 0.35, item, size=12, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════
# Slide 10: Ask / CTA
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, 1, 1.0, 11, 0.5, "NEXT STEPS", size=14, color=MID_GRAY,
         alignment=PP_ALIGN.CENTER)

add_text(slide, 1, 1.8, 11, 1.0, "같이 만들어 갈\n다음 단계",
         size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 3개 CTA 카드
ctas = [
    ("1️⃣  파일럿 도입", "5~10개 제조기업 대상\n3개월 무료 파일럿\n→ 실증 데이터 확보", ACCENT_BLUE),
    ("2️⃣  투자 라운드", "시리즈 A 준비\n글로벌 제조기업 확장\n→ 개발/영업 팀 확장", ACCENT_GREEN),
    ("3️⃣  파트너십", "SAP 파트너 프로그램\nMES/QMS 벤더 협업\n→ 채널 확장", ACCENT_ORANGE),
]

for i, (title, body, color) in enumerate(ctas):
    left = 0.8 + i * 4.1
    add_shape(slide, left, 3.5, 3.8, 2.2)
    add_text(slide, left + 0.3, 3.65, 3.2, 0.4, title, size=18, color=color, bold=True)
    add_text(slide, left + 0.3, 4.15, 3.2, 1.3, body, size=14, color=LIGHT_GRAY)

# 연락처
add_text(slide, 1, 6.2, 11, 0.5,
         "SAP Knowledge Hub  ·  v7.1.0  ·  2026",
         size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
add_text(slide, 1, 6.6, 11, 0.4,
         "감사합니다",
         size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# 저장
# ═══════════════════════════════════════════════
desktop = os.path.expandvars(r"%USERPROFILE%\OneDrive\바탕 화면")
if not os.path.isdir(desktop):
    desktop = os.path.expandvars(r"%USERPROFILE%\Desktop")

output_path = os.path.join(desktop, "SAP_Knowledge_Hub_v7.1_발표.pptx")
prs.save(output_path)
print(f"✅ PPT 생성 완료: {output_path}")
print(f"   슬라이드 수: {len(prs.slides)}")
